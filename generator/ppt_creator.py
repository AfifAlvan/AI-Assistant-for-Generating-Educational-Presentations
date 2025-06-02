from pptx import Presentation
import os

def buat_ppt(judul, nama, isi, output_path="output/hasil_materi.pptx"):
    MAX_SLIDE2_LENGTH = 600  # maksimal karakter untuk slide ringkasan

    if len(isi) > MAX_SLIDE2_LENGTH:
        isi = isi[:MAX_SLIDE2_LENGTH].rstrip() + "..."

    prs = Presentation()

    # Slide 1 - Judul dan pembuat
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # layout judul
    slide1.shapes.title.text = f"Matery: {judul}"
    slide1.placeholders[1].text = f"Created by {nama}"

    # Slide 2 - Ringkasan materi
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # layout judul+konten
    slide2.shapes.title.text = "Content"
    slide2.placeholders[1].text = isi

    # Slide 3 - Penutup
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Thank you"

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    return output_path
