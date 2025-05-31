from pptx import Presentation

def buat_ppt(judul, nama, isi, output_path="output/hasil_materi.pptx"):
    prs = Presentation()

    # Slide 1 - Judul dan pembuat
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide1.shapes.title.text = judul
    slide1.placeholders[1].text = f"Dibuat oleh {nama}"

    # Slide 2 - Deskripsi topik (ringkasan)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Deskripsi Topik"
    slide2.placeholders[1].text = isi

    # Slide 3 - Isi beberapa topik (untuk contoh, pakai isi juga)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Isi Topik"
    slide3.placeholders[1].text = isi  # bisa diganti dengan isi lain jika ada

    # Slide 4 - Terima kasih
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Terima Kasih"
    slide4.placeholders[1].text = "Terima kasih atas perhatian Anda."

    # Simpan file PPT
    prs.save(output_path)
    return output_path
